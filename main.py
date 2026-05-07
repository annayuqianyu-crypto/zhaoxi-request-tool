import json
import re
import io
import os
import secrets
import urllib.request
import urllib.parse
from datetime import datetime, timedelta
from typing import List, Union, Optional

from fastapi import FastAPI, HTTPException, UploadFile, File, Response, Header
from fastapi.staticfiles import StaticFiles
from fastapi.responses import HTMLResponse, StreamingResponse
from pydantic import BaseModel
from openai import OpenAI, AsyncOpenAI
from dotenv import load_dotenv

load_dotenv()

app = FastAPI()

# 默认 AI 配置（公司统一 DeepSeek 账号；用户在页面填写自己的 Key 时可覆盖）
_DEFAULT_API_KEY = os.environ.get("WOLFAI_API_KEY", "sk-603a729e51d54a82bf8b8de3e06530b4")
_DEFAULT_API_URL = os.environ.get("WOLFAI_API_URL", "https://api.deepseek.com")
_DEFAULT_MODEL   = os.environ.get("WOLFAI_MODEL",   "deepseek-v4-pro")

def _get_client(api_key: Optional[str] = None, api_url: Optional[str] = None) -> OpenAI:
    """根据请求头中的 key/url 创建同步 OpenAI 客户端；未提供则使用环境变量默认值。"""
    key = (api_key or "").strip() or _DEFAULT_API_KEY
    url = (api_url or "").strip() or _DEFAULT_API_URL
    if not key:
        raise HTTPException(400, "未配置 API Key，请在页面右上角「⚙ API设置」中填写")
    return OpenAI(api_key=key, base_url=url)

def _get_async_client(api_key: Optional[str] = None, api_url: Optional[str] = None) -> AsyncOpenAI:
    """创建异步 OpenAI 客户端，用于流式响应。"""
    key = (api_key or "").strip() or _DEFAULT_API_KEY
    url = (api_url or "").strip() or _DEFAULT_API_URL
    if not key:
        raise HTTPException(400, "未配置 API Key，请在页面右上角「⚙ API设置」中填写")
    return AsyncOpenAI(api_key=key, base_url=url)

def _get_model(api_model: Optional[str] = None) -> str:
    """返回本次请求使用的模型名；未指定则取环境变量默认值。"""
    return (api_model or "").strip() or _DEFAULT_MODEL

def _truncate_history(history: list, max_pairs: int = 6) -> list:
    """
    截断对话历史，防止超长上下文导致 Render 30s 超时（HTTP 502）。
    策略：保留第 1 条消息（初始背景），然后保留最后 max_pairs 轮（每轮=user+assistant）。
    """
    if not history:
        return history
    max_msgs = max_pairs * 2  # 每轮2条消息
    if len(history) <= max_msgs + 1:
        return history
    # 第0条保留（初始描述背景），其余取最后 max_msgs 条
    return [history[0]] + history[-max_msgs:]


def _clean_mermaid(raw: str) -> str:
    """
    清理 AI 生成的 Mermaid 代码，移除最常见的导致 parse-error 的语法：
    1. 代码围栏（```mermaid … ```）
    2. :::className 行内类名后缀（mermaid v10 解析不稳定）
    3. 独立的 class 应用行（暂留 classDef 定义行，让颜色仍能生效）
    4. 删除第一个 graph/flowchart 关键字之前的所有非代码文字行
    """
    # 移除 ``` 围栏
    cleaned = re.sub(r"```(?:mermaid)?\s*", "", raw).replace("```", "").strip()

    lines = cleaned.split("\n")
    result = []
    found_graph = False
    for line in lines:
        stripped = line.strip()
        # 跳过 graph/flowchart 之前的非代码行
        if not found_graph:
            if re.match(r"^(graph|flowchart)\s+(TD|LR|BT|RL)", stripped, re.IGNORECASE):
                found_graph = True
            else:
                continue
        # 移除行内 :::className（替换为空）
        line = re.sub(r":::[\w-]+", "", line)
        result.append(line)

    return "\n".join(result).strip() if result else cleaned

# ─────────────────────────────────────────────
# Supabase REST API（走 HTTPS，自动通过系统代理）
# 彻底替代 psycopg2 直连，解决 TCP/5432 被代理拦截问题
# ─────────────────────────────────────────────
SUPABASE_URL = os.environ.get("SUPABASE_URL", "https://nbjdukzpjblpavnmmwmm.supabase.co").strip()
SUPABASE_KEY = os.environ.get("SUPABASE_KEY", "").strip()


def _sb(method: str, path: str, *, data=None, params: dict = None, extra_headers: dict = None):
    """
    Supabase PostgREST REST 调用。
    urllib.request 自动读取 HTTPS_PROXY 环境变量，无需额外配置。
    """
    url = SUPABASE_URL + path
    if params:
        parts = [
            f"{urllib.parse.quote(str(k))}="
            f"{urllib.parse.quote(str(v), safe='.,*-_:/()!')}"
            for k, v in params.items()
        ]
        url += "?" + "&".join(parts)

    headers = {
        "apikey": SUPABASE_KEY,
        "Authorization": f"Bearer {SUPABASE_KEY}",
        "Content-Type": "application/json",
        "Prefer": "return=representation",
    }
    if extra_headers:
        headers.update(extra_headers)

    body = json.dumps(data, ensure_ascii=False).encode("utf-8") if data is not None else None
    req = urllib.request.Request(url, data=body, headers=headers, method=method)

    try:
        with urllib.request.urlopen(req, timeout=15) as resp:
            raw = resp.read()
            return json.loads(raw) if raw else {}
    except urllib.error.HTTPError as e:
        err = e.read().decode(errors="replace")
        print(f"[Supabase HTTPError] {method} {path} → {e.code}: {err[:300]}")
        raise HTTPException(e.code, f"数据库错误 {e.code}: {err[:300]}")
    except urllib.error.URLError as e:
        print(f"[Supabase URLError] {method} {path} → {e.reason}")
        raise HTTPException(503, f"数据库连接失败: {e.reason}")
    except Exception as e:
        print(f"[Supabase Exception] {method} {path} → {type(e).__name__}: {e}")
        raise HTTPException(500, f"数据库异常: {type(e).__name__}: {str(e)[:200]}")


def _sb_count(path: str, params: dict = None) -> int:
    """返回符合条件的行数（select=id 然后在 Python 侧计数）。"""
    try:
        rows = _sb("GET", path, params={**(params or {}), "select": "id"})
        return len(rows) if isinstance(rows, list) else 0
    except Exception:
        return 0


def init_db():
    """
    验证 Supabase 表已存在并可访问。
    若表不存在，请在 Supabase 控制台 SQL Editor 中执行：
      CREATE TABLE IF NOT EXISTS users (
          id TEXT PRIMARY KEY, email TEXT UNIQUE NOT NULL,
          name TEXT NOT NULL, token TEXT UNIQUE NOT NULL,
          is_admin INTEGER DEFAULT 0,
          created_at TEXT NOT NULL, last_active TEXT);
      CREATE TABLE IF NOT EXISTS sessions (
          id TEXT PRIMARY KEY, user_id TEXT NOT NULL,
          title TEXT, history TEXT, requirements TEXT, mermaid TEXT,
          completeness INTEGER DEFAULT 0,
          created_at TEXT NOT NULL, updated_at TEXT NOT NULL);
    """
    _sb("GET", "/rest/v1/users",    params={"select": "id", "limit": "1"})
    _sb("GET", "/rest/v1/sessions", params={"select": "id", "limit": "1"})


try:
    init_db()
    print("[OK] Supabase REST API connected and tables verified.")
except Exception as _e:
    print(f"[WARN] Supabase check: {_e}")
    print("[INFO] If tables are missing, run the CREATE TABLE SQL in Supabase dashboard.")


def require_token(authorization: Optional[str] = Header(None)):
    """从 Authorization: Bearer <token> 中验证用户"""
    if not authorization or not authorization.startswith("Bearer "):
        raise HTTPException(401, "未登录，请先登录")
    token = authorization[7:]
    rows = _sb("GET", "/rest/v1/users", params={"token": f"eq.{token}", "select": "*"})
    if not rows or not isinstance(rows, list):
        raise HTTPException(401, "登录已失效，请重新登录")
    return rows[0]


def require_admin(authorization: Optional[str] = Header(None)):
    user = require_token(authorization)
    if not user.get("is_admin"):
        raise HTTPException(403, "无管理员权限")
    return user


# ─────────────────────────────────────────────
# System prompt：Grill-Me 风格 + 结构化需求输出
# ─────────────────────────────────────────────
SYSTEM_PROMPT = """你是朝曦金融机构高端客户增值服务平台的资深业务分析师，精通公司五大核心服务体系：
1. 家族财富架构服务（境内外信托架构、资本市场投融资、全球资产配置、股权激励）
2. 全球税务规划（美/加/新加坡税务合规、上市公司股东减持税务、跨境架构涉税、CRS/FATCA）
3. 全球法律咨询（资本金融、公司商事、婚姻家事、企业合规、不良资产处置）
4. 资本市场服务（上市公司股东减持、股东融资退出、拟上市公司架构、私募基金架构）
5. 企业治理（ESG合规、组织效能提升、股权激励体系、接班人规划）

五大客群：上市公司股东、新兴/传统企业家、国际化人士（美/加/澳/英）、企业家太太、高净值家族。

任务：通过深度访谈，将朝曦业务团队模糊的想法转化为清晰、可落地的IT系统需求。

## 访谈方法论

**第一步：场景剧本破冰**
不要直接问"你要什么功能"，请业务方讲一个真实工作故事：
"请描述最近您或同事实际遇到的一个具体情境，就像在给我讲故事一样。"

**第二步：结构化深挖（每次只问一个最关键的问题）**
根据故事，判断哪个维度最不清晰，逐一追问：
- 【角色与权限】谁在用？什么职级？权限边界？
- 【触发条件】什么情况下启动流程？频率？
- 【输入数据】需要哪些信息？从哪里来？格式？
- 【处理流程】每个步骤是什么？谁来操作？
- 【审批链路】谁审批？几级？时限要求？
- 【输出结果】产出是什么？发给谁？存在哪里？
- 【异常处理】出错怎么办？有哪些特殊情况？
- 【合规要求】有哪些监管规定或内部制度必须遵守？

**领域专项追问（自动识别业务领域）：**
- 税务场景：申报期限、税种、纳税主体、税务机关对接、凭证留存年限
- 法律场景：合同类型、审阅层级、印章管控、律师介入时机、诉讼风险分级
- 资本市场：交易品种、风控限额、监管报送、估值方式、结算周期

## 追问原则
1. 每次只问一个问题，绝不连问
2. 不接受模糊回答：听到"大概""差不多"时追问具体数字或细节
3. completeness 达 80 以上才建议进入总结
4. 全程使用中文

## 严格输出格式（只输出 JSON，不要任何其他文字）
{
  "message": "AI回应文字 + 下一个追问问题",
  "stage": "opening 或 exploring 或 drilling 或 summarizing",
  "completeness": 0到100整数
}"""

OPENING_MESSAGE = {
    "message": "您好！我是朝曦的业务需求分析助手，熟悉家族财富架构、全球税务规划、资本市场服务、法律咨询和企业治理五大业务体系。\n\n**我的工作方式：** 通过深度访谈把模糊想法变成清晰的IT需求，对话结束后点击右侧各板块的「⚡ 生成」按钮，可一键生成需求梳理、流程图、线框图和PRD文档，最终可直接交付IT团队。\n\n**请按以下框架描述您的场景**（能说多少说多少，其余我来追问）：\n\n🏢 **所在部门**：哪个业务条线或岗位会使用这个系统？\n　　（如：税务团队 / 架构师团队 / 资本市场组 / 客服中台）\n😣 **业务痛点**：目前这件事最大的困难或效率瓶颈是什么？\n🎯 **期望解决**：您希望系统能帮您做到什么，达到什么效果？\n📥 **涉及输入**：需要录入或上传哪些信息？数据从哪里来？\n📤 **期望输出**：系统最终要产出什么？发给谁？存在哪里？\n\n您也可以直接上传相关文件（Word/Excel/PPT）或使用 🎤 语音描述，我会帮您提炼需求。",
    "stage": "opening",
    "completeness": 0
}


# ─────────────────────────────────────────────
# Data models
# ─────────────────────────────────────────────
class Message(BaseModel):
    role: str
    content: str

class LoginRequest(BaseModel):
    name: str
    email: str

class SaveSessionRequest(BaseModel):
    session_id: Optional[str] = None
    title: str
    history: list
    requirements: list
    mermaid: str = ""
    wireframe_mermaid: str = ""
    ui_wireframe_mermaid: str = ""
    prd_content: str = ""
    completeness: int = 0

class ChatRequest(BaseModel):
    history: List[Message]
    message: str


class ExportRequest(BaseModel):
    history: List[Message]
    requirements: List[Union[dict, str]]
    mermaid: str


class WordExportReq(BaseModel):
    type: str  # "requirements" or "prd"
    requirements: List[Union[dict, str]] = []
    prd_text: str = ""


class GenerateRequest(BaseModel):
    history: List[Message]
    requirements: List[Union[dict, str]]
    type: str  # summary | wireframe | preview | flowchart

class DemoRequest(BaseModel):
    history: List[Message]
    requirements: List[Union[dict, str]] = []
    mermaid: str = ""


# ─────────────────────────────────────────────
# File extraction helpers
# ─────────────────────────────────────────────
def extract_docx(content: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(content))
    texts = []
    for para in doc.paragraphs:
        if para.text.strip():
            texts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
            if row_text:
                texts.append(row_text)
    return "\n".join(texts)


def extract_xlsx(content: bytes) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(content), data_only=True)
    texts = []
    for sheet in wb.worksheets:
        texts.append(f"[工作表：{sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            vals = [str(v) for v in row if v is not None and str(v).strip()]
            if vals:
                texts.append(" | ".join(vals))
    return "\n".join(texts[:300])


def extract_pptx(content: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(content))
    texts = []
    for i, slide in enumerate(prs.slides, 1):
        texts.append(f"\n[幻灯片 {i}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        texts.append(para.text)
    return "\n".join(texts)


def parse_ai_response(raw: str) -> dict:
    cleaned = re.sub(r"```(?:json)?\s*", "", raw).replace("```", "").strip()

    # 第一次：直接解析
    try:
        return json.loads(cleaned)
    except Exception:
        pass

    # 第二次：修复 JSON 字符串值内部的裸换行符（AI 有时不转义换行）
    # 用状态机逐字符扫描，只在字符串内部转换 \n / \r
    try:
        out, in_str, esc = [], False, False
        for ch in cleaned:
            if esc:
                out.append(ch); esc = False
            elif ch == '\\':
                out.append(ch); esc = True
            elif ch == '"':
                out.append(ch); in_str = not in_str
            elif in_str and ch == '\n':
                out.append('\\n')
            elif in_str and ch == '\r':
                out.append('\\r')
            else:
                out.append(ch)
        return json.loads(''.join(out))
    except Exception:
        pass

    # 第三次：正则提取最外层 JSON 对象
    match = re.search(r"\{[\s\S]*\}", cleaned)
    if match:
        try:
            return json.loads(match.group())
        except Exception:
            pass

    # 最终 fallback：提取 message 字段文本
    # 用 json.loads 解码（正确处理 \uXXXX 等转义），避免 .decode('unicode_escape') 导致中文乱码
    msg_match = re.search(r'"message"\s*:\s*"((?:[^"\\]|\\.)*)"', cleaned)
    if msg_match:
        try:
            fallback_msg = json.loads('"' + msg_match.group(1) + '"')
        except Exception:
            fallback_msg = msg_match.group(1)
    else:
        fallback_msg = "AI 正在分析中，请稍候…"
    return {
        "message": fallback_msg,
        "stage": "exploring",
        "completeness": 10,
    }


# ─────────────────────────────────────────────
# API routes
# ─────────────────────────────────────────────
@app.get("/api/init")
async def init():
    return OPENING_MESSAGE


@app.post("/api/test-connection")
async def test_connection(x_api_key: Optional[str] = Header(None),
                          x_api_url: Optional[str] = Header(None),
                          x_api_model: Optional[str] = Header(None)):
    """用极小的 completion 请求验证 API Key / URL / Model 是否有效"""
    client = _get_client(x_api_key, x_api_url)
    model  = _get_model(x_api_model)
    try:
        client.chat.completions.create(
            model=model,
            messages=[{"role": "user", "content": "hi"}],
            max_tokens=1,
        )
        return {"ok": True, "model": model}
    except Exception as e:
        raise HTTPException(400, f"连接失败：{str(e)[:200]}")


@app.post("/api/upload")
async def upload_file(file: UploadFile = File(...)):
    filename = file.filename or ""
    content = await file.read()
    ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""

    try:
        if ext == "docx":
            text = extract_docx(content)
        elif ext == "xlsx":
            text = extract_xlsx(content)
        elif ext in ("pptx", "ppt"):
            text = extract_pptx(content)
        else:
            raise HTTPException(400, "仅支持 .docx、.xlsx、.pptx 格式")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"文件解析失败：{e}")

    MAX = 8000
    if len(text) > MAX:
        text = text[:MAX] + f"\n…（内容过长，已截取前 {MAX} 字）"

    return {"filename": filename, "extracted_text": text, "char_count": len(text)}


@app.post("/api/chat")
async def chat(req: ChatRequest,
               x_api_key: Optional[str] = Header(None),
               x_api_url: Optional[str] = Header(None),
               x_api_model: Optional[str] = Header(None)):
    """
    流式对话接口：逐 token 推送 AI 输出，彻底绕开 Render 30s HTTP 超时限制。
    前端收到 data: [DONE] 后将累积文本解析为 JSON。
    支持永续对话，不限轮次。
    """
    try:
        client = _get_async_client(x_api_key, x_api_url)
    except HTTPException:
        raise
    model = _get_model(x_api_model)

    # 保留第0条（初始背景）+ 最近4轮，控制 token 量
    truncated = _truncate_history(list(req.history), max_pairs=4)
    messages  = [{"role": m.role, "content": m.content} for m in truncated]
    messages.append({"role": "user", "content": req.message})

    async def generate():
        try:
            stream = await client.chat.completions.create(
                model=model,
                max_tokens=2048,
                messages=[{"role": "system", "content": SYSTEM_PROMPT}] + messages,
                stream=True,
            )
            async for chunk in stream:
                delta = chunk.choices[0].delta.content or ""
                if delta:
                    # 逐块推送原始 token，前端累积后统一解析
                    yield delta.encode("utf-8")
            # 流结束标志
            yield b"\n\ndata: [DONE]"
        except Exception as e:
            # 将错误嵌入流末尾，前端检测后显示友好报错
            err = json.dumps({"__stream_error__": str(e)}, ensure_ascii=False)
            yield f"\n\ndata: [ERR]{err}".encode("utf-8")

    return StreamingResponse(
        generate(),
        media_type="text/plain; charset=utf-8",
        headers={"X-Accel-Buffering": "no"},  # 禁用 Nginx/Render 的响应缓冲，确保实时推送
    )


@app.post("/api/export")
async def export_doc(req: ExportRequest,
                     x_api_key: Optional[str] = Header(None),
                     x_api_url: Optional[str] = Header(None),
                     x_api_model: Optional[str] = Header(None)):
    client = _get_client(x_api_key, x_api_url)
    model  = _get_model(x_api_model)
    history_text = "\n\n".join(
        f"**{'业务方' if m.role == 'user' else '分析师'}**：{m.content}"
        for m in req.history
    )

    req_lines = []
    for r in req.requirements:
        if isinstance(r, dict):
            req_lines.append(f"\n### {r.get('module','')} · {r.get('feature','')}")
            req_lines.append(f"**功能说明**：{r.get('description','')}")
            if r.get("wireframe"):
                req_lines.append(f"**界面草图**：{r['wireframe']}")
            if r.get("process"):
                req_lines.append("**处理流程**：" + " → ".join(r["process"]))
            if r.get("inputs"):
                req_lines.append("**输入数据**：" + "、".join(r["inputs"]))
            if r.get("outputs"):
                req_lines.append("**输出结果**：" + "、".join(r["outputs"]))
        else:
            req_lines.append(f"- {r}")
    req_text = "\n".join(req_lines) or "（请基于访谈记录梳理需求）"

    prompt = f"""根据以下访谈记录和已梳理需求，生成一份专业的IT需求规格说明书（Markdown格式）。

## 访谈记录
{history_text}

## 已梳理需求
{req_text}

请输出包含以下章节的完整文档：
1. **项目背景与目标**
2. **涉及角色与权限**
3. **业务流程说明**
4. **功能需求清单**（每条需求含描述、验收标准）
5. **界面设计要点**（关键页面布局说明）
6. **数据需求**（关键字段、来源、格式、存储要求）
7. **非功能性需求**（性能、安全、合规条款）
8. **待确认事项**（仍需业务方明确的问题）
9. **验收标准**

要求：专业严谨，结构清晰，可直接交付IT团队进行工作量评估。"""

    try:
        response = client.chat.completions.create(
            model=model,
            max_tokens=4096,
            messages=[{"role": "user", "content": prompt}],
        )
    except Exception as e:
        raise HTTPException(500, f"导出失败：{e}")

    return {"document": response.choices[0].message.content}


@app.post("/api/export/code")
async def export_frontend_code(req: ExportRequest,
                                x_api_key: Optional[str] = Header(None),
                                x_api_url: Optional[str] = Header(None),
                                x_api_model: Optional[str] = Header(None)):
    """流式生成各功能页面的 HTML/CSS/JS 实现代码，供 IT 开发人员直接参考"""
    client_async = _get_async_client(x_api_key, x_api_url)
    model = _get_model(x_api_model)

    gen_history = _truncate_history(list(req.history), max_pairs=10)
    history_text = "\n\n".join(
        f"{'业务方' if m.role == 'user' else '分析师'}：{m.content}"
        for m in gen_history
    )
    req_text = "\n".join(
        f"【{r.get('module','')} · {r.get('feature','')}】{r.get('description','')}"
        if isinstance(r, dict) else f"- {r}"
        for r in req.requirements
    ) or "（根据访谈记录梳理需求）"

    prompt = f"""根据以下访谈记录和需求清单，为每个功能模块生成对应的前端实现代码。

## 访谈记录
{history_text}

## 需求清单
{req_text}

## 输出要求
1. 输出一个完整的、可直接在浏览器打开的 HTML 文件
2. 技术栈：原生 HTML5 + CSS3 + JavaScript（不使用外部框架/CDN）
3. 文件结构：
   - 顶部：目录索引（锚点链接，点击可跳转到各功能页面）
   - 每个功能模块用 <section id="feature-N"> 包裹，标注模块名和功能名
   - 各模块间有分隔线
4. 每个功能模块的代码必须包含：
   - 页面标题 + 面包屑导航
   - 若有查询/筛选：生成 <input>/<select> 筛选控件 + 查询按钮
   - 若有列表/表格：生成 <table> 含表头和2-3行示例数据
   - 若有表单：生成 <form> 含对应 <input>/<select>/<textarea> 字段
   - 操作按钮：提交/保存/取消（带 JS 确认弹窗）
   - 简单的表单验证 JS（必填项检查）
5. CSS 风格：简洁现代的企业内网管理系统，使用以下配色：
   - 主色：#1a365d（深蓝导航）
   - 辅色：#3182ce（按钮/链接）
   - 背景：#f7f8fa
   - 卡片/表单背景：#ffffff，border-radius:8px，box-shadow 轻阴影
   - 表格：奇偶行交替色 #f9fafb / #ffffff
6. 只输出完整 HTML 代码，从 <!DOCTYPE html> 开始，不要任何代码块标记或说明文字"""

    async def stream_code():
        try:
            stream = await client_async.chat.completions.create(
                model=model, max_tokens=4096,
                messages=[{"role": "user", "content": prompt}],
                stream=True,
            )
            async for chunk in stream:
                delta = chunk.choices[0].delta.content or ""
                if delta:
                    yield delta.encode("utf-8")
            yield b"\n\ndata: [DONE]"
        except Exception as e:
            err = json.dumps({"__stream_error__": str(e)}, ensure_ascii=False)
            yield f"\n\ndata: [ERR]{err}".encode("utf-8")

    return StreamingResponse(
        stream_code(),
        media_type="text/plain; charset=utf-8",
        headers={"X-Accel-Buffering": "no"},
    )


@app.post("/api/generate")
async def generate_section(req: GenerateRequest,
                            x_api_key: Optional[str] = Header(None),
                            x_api_url: Optional[str] = Header(None),
                            x_api_model: Optional[str] = Header(None)):
    client = _get_client(x_api_key, x_api_url)
    model  = _get_model(x_api_model)
    # 截断历史，最多取最后 10 轮（避免生成接口也超时）
    gen_history = _truncate_history(list(req.history), max_pairs=10)
    history_text = "\n\n".join(
        f"{'业务方' if m.role == 'user' else '分析师'}：{m.content}"
        for m in gen_history
    )

    if req.type == "flowchart":
        prompt = f"""根据以下访谈记录，生成一份完整的 Mermaid 业务流程图。

访谈记录：
{history_text}

【严格遵守以下 Mermaid 语法规则，违反任一条均会导致图表无法渲染】：
1. 第一行必须是：graph TD
2. 节点 ID 只允许使用英文字母和数字（如 A1、stepB、end1），严禁使用中文、空格或特殊符号作为 ID
3. 节点标签必须用双引号包裹：A1["中文描述"]；判断节点用花括号：D1{{"是否满足？"}}；开始/结束用圆括号：S(["开始"])
4. 连线用 --> 或 -->|"说明文字"|
5. 可以用 subgraph 对阶段分组：subgraph "阶段名称" ... end
6. 如需节点样式，使用 classDef + class 语句（不要用 :::className 内联写法）：
   classDef start fill:#d1fae5,stroke:#10b981,color:#065f46
   classDef proc fill:#dbeafe,stroke:#3b82f6,color:#1e40af
   classDef dec fill:#fce7f3,stroke:#ec4899,color:#831843
   class S,E start
   class A1,B2 proc
7. 标签内不能出现未转义的英文双引号，用中文标点替代
8. 只输出 Mermaid 代码，不要任何说明、注释或多余文字"""
        try:
            response = client.chat.completions.create(
                model=model, max_tokens=2000,
                messages=[{"role": "user", "content": prompt}]
            )
        except Exception as e:
            raise HTTPException(500, f"生成失败：{e}")
        raw = response.choices[0].message.content
        cleaned = _clean_mermaid(raw)
        return {"type": "flowchart", "mermaid": cleaned, "requirements": []}

    elif req.type == "wireframe":
        prompt = f"""根据以下访谈记录，生成一份供 IT 人员使用的功能流程设计图（Mermaid 格式）。

访谈记录：
{history_text}

【严格遵守以下 Mermaid 语法规则，违反任一条均会导致图表无法渲染】：
1. 第一行必须是：graph TD
2. 节点 ID 只允许使用英文字母和数字（如 U1、S2、D3），严禁使用中文、空格或特殊符号作为 ID
3. 节点标签必须用双引号包裹：
   - 操作步骤（矩形）：U1["用户提交申请"]
   - 判断节点（菱形）：D1{{"审核是否通过？"}}
   - 开始/结束（圆角）：S(["开始"]) 或 E(["结束"])
4. 连线用 --> 或 -->|"说明"|
5. 多方案/多路径用 subgraph 分组：subgraph "方案A" ... end
6. 使用 classDef + class 语句为节点着色（不要用 :::className 内联写法）：
   classDef user fill:#fff3cd,stroke:#f59e0b,color:#92400e
   classDef sys fill:#dbeafe,stroke:#3b82f6,color:#1e40af
   classDef dec fill:#fce7f3,stroke:#ec4899,color:#831843
   classDef ep fill:#d1fae5,stroke:#10b981,color:#065f46
   class U1,U2 user
   class S1,S2 sys
   class D1,D2 dec
   class START,END ep
7. 标签内不能出现未转义的英文双引号，用中文标点替代
8. 只输出 Mermaid 代码，不要任何说明、注释或多余文字"""
        try:
            response = client.chat.completions.create(
                model=model, max_tokens=2500,
                messages=[{"role": "user", "content": prompt}]
            )
        except Exception as e:
            raise HTTPException(500, f"生成失败：{e}")
        raw = response.choices[0].message.content
        cleaned = _clean_mermaid(raw)
        return {"type": "wireframe", "mermaid": cleaned, "requirements": []}

    elif req.type in ("summary", "preview"):
        if req.type == "summary":
            extra = "只需要 module/feature/description/process/inputs/outputs 字段，无需 wireframe。"
        else:
            extra = "必须包含 wireframe 字段，格式：顶部：...\\n中部：...\\n底部：..."

        prompt = f"""根据以下访谈记录，提取并整理所有已明确的功能需求，以 JSON 数组格式输出。

访谈记录：
{history_text}

要求：
- 每个需求对象包含：module（模块名）、feature（功能名）、description（功能说明）、wireframe（界面布局，格式"顶部：...\\n中部：...\\n底部：..."）、process（操作步骤数组）、inputs（输入项数组）、outputs（输出项数组）
- {extra}
- 只输出 JSON 数组，不要任何其他文字
- 确保 JSON 格式合法"""
        try:
            response = client.chat.completions.create(
                model=model, max_tokens=3000,
                messages=[{"role": "user", "content": prompt}]
            )
        except Exception as e:
            raise HTTPException(500, f"生成失败：{e}")
        raw = response.choices[0].message.content
        cleaned = re.sub(r"```(?:json)?\s*", "", raw).replace("```", "").strip()
        try:
            data = json.loads(cleaned)
            if isinstance(data, list):
                return {"type": req.type, "requirements": data, "mermaid": ""}
        except Exception:
            match = re.search(r"\[[\s\S]*\]", cleaned)
            if match:
                try:
                    return {"type": req.type, "requirements": json.loads(match.group()), "mermaid": ""}
                except Exception:
                    pass
        return {"type": req.type, "requirements": [], "mermaid": ""}

    elif req.type == "ui_wireframe":
        prompt = f"""你是一位资深的B端产品架构师与UX设计师。根据以下访谈记录，生成一份面向IT交付的系统线框图。

访谈记录：
{history_text}

## 思考步骤（内部思考，不输出过程）
1. 解构需求：识别3-5个核心业务流关键词
2. 映射模块：将关键词映射为系统功能模块
3. 定义动作与产出：每模块列出关键业务动作和对应产出物
4. 组装线框图：按标准B端布局组装完整设计

## 输出格式（必须严格按此格式，不能偏离）

[FLOW]
graph TD
（此处输出业务流程Mermaid代码，要求如下：）
（1）节点 ID 只用英文字母数字，如 U1 S2 D1 OUT1
（2）节点标签必须用双引号，中文描述，内部不得出现英文双引号
（3）用不同节点形状区分：用户动作用矩形 U1["动作"]，系统处理用圆角 S1(["处理"])，判断用菱形 D1{{"判断？"}}，产出物用六边形 O1{{"产出"}}
（4）连线：U1 --> S1 或 S1 -->|"条件"| D1
（5）用 subgraph 对主要业务阶段分组
（6）节点样式：白底黑字黑边，classDef user fill:#ffffff,stroke:#000000,color:#000000 / classDef sys fill:#f5f5f5,stroke:#000000,color:#000000 / classDef output fill:#fffde7,stroke:#000000,color:#000000
（7）最后一行必须是：linkStyle default stroke:#FFD700,stroke-width:3px
[/FLOW]

[SPEC]
（此处输出核心页面规格说明，必须包含以下各项）

**系统名称**：[根据需求命名]
**业务流程总结**：[用5-8个关键词描述核心业务流，如：客户录入 → 需求分析 → 方案匹配 → 审批交付]

---

**核心页面：[主要工作台页面名称]**

布局结构：
- 顶部导航栏：[系统名称] | 🔍 全局搜索框 | 通知铃 | 👤 [当前用户名/租户名]
- 左侧菜单（一级/二级）：
  ├── [模块1图标] [模块1名称]
  │   ├── [子菜单1]（当前激活）
  │   └── [子菜单2]
  └── [模块2图标] [模块2名称]
- 主体-操作区：[+ 新建] [筛选▼ 状态/时间/类型] [批量操作▼] [导出] | 右侧：搜索框
- 主体-数据区（表格）：
  表头字段：序号 | [字段1] | [字段2] | [字段3] | [字段4] | 状态 | 操作
  示例数据行1：001 | [示例值] | [示例值] | [示例值] | [示例值] | ● 处理中 | 查看 编辑
  示例数据行2：002 | [示例值] | [示例值] | [示例值] | [示例值] | ✓ 已完成 | 查看 编辑
- 主体-右侧详情抽屉（点击行后展开）：
  Tab1：基本信息 / Tab2：[业务相关信息] / Tab3：操作日志
  底部操作按钮：[主操作按钮] [次要操作] [取消]
- 空状态引导页：图标 + "暂无数据，点击[+ 新建]开始创建" + 按钮
- 加载态：骨架屏占位 + 进度提示

关键交互流：
├── 点击[+ 新建] → 唤起[新建表单弹窗/抽屉] → 产生产出[新记录保存，刷新列表]
├── 点击[某行数据] → 唤起[右侧详情抽屉展开] → 产生产出[完整信息展示]
├── 点击[筛选▼] → 唤起[筛选条件面板] → 产生产出[过滤后的数据列表]
├── 点击[主操作按钮] → 唤起[二次确认弹窗] → 产生产出[状态变更，操作日志记录]
└── 点击[导出] → 唤起[导出配置选择] → 产生产出[Excel/PDF文件下载]

字段与实体映射：
| 实体名 | 字段名 | 数据类型 | 前端组件类型 | 说明 |
|--------|--------|----------|------------|------|
| [实体1] | [字段1] | String | 单行输入框 | [必填/选填] |
| [实体1] | [字段2] | Enum | 下拉选择器 | 枚举值：[值1/值2/值3] |
| [实体1] | [字段3] | Date | 日期选择器 | 格式：YYYY-MM-DD |
| [实体1] | [字段4] | Number | 数字输入框 | [单位/范围说明] |
| [实体1] | [字段5] | Text | 多行文本域 | 最大500字 |
| [实体2] | [字段1] | String | 只读文本 | 系统自动生成 |

---

（如有多个核心页面，继续按上述格式列出第二、第三个页面）
[/SPEC]

重要提示：
- [FLOW]和[/FLOW]之间只输出合法Mermaid代码，不含任何说明文字
- [SPEC]和[/SPEC]之间输出结构化规格文本，用实际业务内容替换所有[括号内的示例占位符]
- linkStyle default 必须放在Mermaid代码的最最后一行
- 所有节点ID仅用英文字母数字，不含中文、空格、特殊字符"""
        try:
            response = client.chat.completions.create(
                model=model, max_tokens=2500,
                messages=[{"role": "user", "content": prompt}]
            )
        except Exception as e:
            raise HTTPException(500, f"生成失败：{e}")
        raw = response.choices[0].message.content
        cleaned = _clean_mermaid(raw)
        return {"type": "ui_wireframe", "mermaid": cleaned, "requirements": []}

    raise HTTPException(400, "未知生成类型")


@app.post("/api/generate-stream")
async def generate_section_stream(req: GenerateRequest,
                                   x_api_key: Optional[str] = Header(None),
                                   x_api_url: Optional[str] = Header(None),
                                   x_api_model: Optional[str] = Header(None)):
    """流式生成 Mermaid 图表代码（wireframe / flowchart），逐 token 推送到前端。
    前端可实时显示代码文本，收到 [DONE] 后再触发渲染。"""
    if req.type not in ("wireframe", "flowchart", "ui_wireframe"):
        raise HTTPException(400, "此接口仅支持 wireframe、flowchart 和 ui_wireframe 类型")

    client_async = _get_async_client(x_api_key, x_api_url)
    model = _get_model(x_api_model)
    gen_history = _truncate_history(list(req.history), max_pairs=10)
    history_text = "\n\n".join(
        f"{'业务方' if m.role == 'user' else '分析师'}：{m.content}"
        for m in gen_history
    )

    if req.type == "wireframe":
        prompt = f"""根据以下访谈记录，生成一份供 IT 人员使用的功能流程设计图（Mermaid 格式）。

访谈记录：
{history_text}

【严格遵守以下 Mermaid 语法规则，违反任一条均会导致图表无法渲染】：
1. 第一行必须是：graph TD
2. 节点 ID 只允许使用英文字母和数字（如 U1、S2、D3），严禁使用中文、空格或特殊符号作为 ID
3. 节点标签必须用双引号包裹：
   - 操作步骤（矩形）：U1["用户提交申请"]
   - 判断节点（菱形）：D1{{"审核是否通过？"}}
   - 开始/结束（圆角）：S(["开始"]) 或 E(["结束"])
4. 连线用 --> 或 -->|"说明"|
5. 多方案/多路径用 subgraph 分组：subgraph "方案A" ... end
6. 使用 classDef + class 语句为节点着色（不要用 :::className 内联写法）：
   classDef user fill:#fff3cd,stroke:#f59e0b,color:#92400e
   classDef sys fill:#dbeafe,stroke:#3b82f6,color:#1e40af
   classDef dec fill:#fce7f3,stroke:#ec4899,color:#831843
   classDef ep fill:#d1fae5,stroke:#10b981,color:#065f46
   class U1,U2 user
   class S1,S2 sys
   class D1,D2 dec
   class START,END ep
7. 标签内不能出现未转义的英文双引号，用中文标点替代
8. 只输出 Mermaid 代码，不要任何说明、注释或多余文字"""
    elif req.type == "ui_wireframe":
        prompt = f"""你是一位资深的B端产品架构师与UX设计师。根据以下访谈记录，生成一份面向IT交付的系统线框图。

访谈记录：
{history_text}

## 思考步骤（内部思考，不输出过程）
1. 解构需求：识别3-5个核心业务流关键词
2. 映射模块：将关键词映射为系统功能模块
3. 定义动作与产出：每模块列出关键业务动作和对应产出物
4. 组装线框图：按标准B端布局组装完整设计

## 输出格式（必须严格按此格式，不能偏离）

[FLOW]
graph TD
（此处输出业务流程Mermaid代码，要求如下：）
（1）节点 ID 只用英文字母数字，如 U1 S2 D1 OUT1
（2）节点标签必须用双引号，中文描述，内部不得出现英文双引号
（3）用不同节点形状区分：用户动作用矩形 U1["动作"]，系统处理用圆角 S1(["处理"])，判断用菱形 D1{{"判断？"}}，产出物用六边形 O1{{"产出"}}
（4）连线：U1 --> S1 或 S1 -->|"条件"| D1
（5）用 subgraph 对主要业务阶段分组
（6）节点样式：白底黑字黑边，classDef user fill:#ffffff,stroke:#000000,color:#000000 / classDef sys fill:#f5f5f5,stroke:#000000,color:#000000 / classDef output fill:#fffde7,stroke:#000000,color:#000000
（7）最后一行必须是：linkStyle default stroke:#FFD700,stroke-width:3px
[/FLOW]

[SPEC]
（此处输出核心页面规格说明，必须包含以下各项）

**系统名称**：[根据需求命名]
**业务流程总结**：[用5-8个关键词描述核心业务流，如：客户录入 → 需求分析 → 方案匹配 → 审批交付]

---

**核心页面：[主要工作台页面名称]**

布局结构：
- 顶部导航栏：[系统名称] | 🔍 全局搜索框 | 通知铃 | 👤 [当前用户名/租户名]
- 左侧菜单（一级/二级）：
  ├── [模块1图标] [模块1名称]
  │   ├── [子菜单1]（当前激活）
  │   └── [子菜单2]
  └── [模块2图标] [模块2名称]
- 主体-操作区：[+ 新建] [筛选▼ 状态/时间/类型] [批量操作▼] [导出] | 右侧：搜索框
- 主体-数据区（表格）：
  表头字段：序号 | [字段1] | [字段2] | [字段3] | [字段4] | 状态 | 操作
  示例数据行1：001 | [示例值] | [示例值] | [示例值] | [示例值] | ● 处理中 | 查看 编辑
  示例数据行2：002 | [示例值] | [示例值] | [示例值] | [示例值] | ✓ 已完成 | 查看 编辑
- 主体-右侧详情抽屉（点击行后展开）：
  Tab1：基本信息 / Tab2：[业务相关信息] / Tab3：操作日志
  底部操作按钮：[主操作按钮] [次要操作] [取消]
- 空状态引导页：图标 + "暂无数据，点击[+ 新建]开始创建" + 按钮
- 加载态：骨架屏占位 + 进度提示

关键交互流：
├── 点击[+ 新建] → 唤起[新建表单弹窗/抽屉] → 产生产出[新记录保存，刷新列表]
├── 点击[某行数据] → 唤起[右侧详情抽屉展开] → 产生产出[完整信息展示]
├── 点击[筛选▼] → 唤起[筛选条件面板] → 产生产出[过滤后的数据列表]
├── 点击[主操作按钮] → 唤起[二次确认弹窗] → 产生产出[状态变更，操作日志记录]
└── 点击[导出] → 唤起[导出配置选择] → 产生产出[Excel/PDF文件下载]

字段与实体映射：
| 实体名 | 字段名 | 数据类型 | 前端组件类型 | 说明 |
|--------|--------|----------|------------|------|
| [实体1] | [字段1] | String | 单行输入框 | [必填/选填] |
| [实体1] | [字段2] | Enum | 下拉选择器 | 枚举值：[值1/值2/值3] |
| [实体1] | [字段3] | Date | 日期选择器 | 格式：YYYY-MM-DD |
| [实体1] | [字段4] | Number | 数字输入框 | [单位/范围说明] |
| [实体1] | [字段5] | Text | 多行文本域 | 最大500字 |
| [实体2] | [字段1] | String | 只读文本 | 系统自动生成 |

---

（如有多个核心页面，继续按上述格式列出第二、第三个页面）
[/SPEC]

重要提示：
- [FLOW]和[/FLOW]之间只输出合法Mermaid代码，不含任何说明文字
- [SPEC]和[/SPEC]之间输出结构化规格文本，用实际业务内容替换所有[括号内的示例占位符]
- linkStyle default 必须放在Mermaid代码的最最后一行
- 所有节点ID仅用英文字母数字，不含中文、空格、特殊字符"""
    else:
        prompt = f"""根据以下访谈记录，生成一份完整的 Mermaid 业务流程图。

访谈记录：
{history_text}

【严格遵守以下 Mermaid 语法规则，违反任一条均会导致图表无法渲染】：
1. 第一行必须是：graph TD
2. 节点 ID 只允许使用英文字母和数字（如 A1、stepB、end1），严禁使用中文、空格或特殊符号作为 ID
3. 节点标签必须用双引号包裹：A1["中文描述"]；判断节点用花括号：D1{{"是否满足？"}}；开始/结束用圆括号：S(["开始"])
4. 连线用 --> 或 -->|"说明文字"|
5. 可以用 subgraph 对阶段分组：subgraph "阶段名称" ... end
6. 如需节点样式，使用 classDef + class 语句（不要用 :::className 内联写法）
7. 标签内不能出现未转义的英文双引号，用中文标点替代
8. 只输出 Mermaid 代码，不要任何说明、注释或多余文字"""

    async def stream_gen():
        try:
            max_tok = 4000 if req.type == "ui_wireframe" else 2500
            stream = await client_async.chat.completions.create(
                model=model, max_tokens=max_tok,
                messages=[{"role": "user", "content": prompt}],
                stream=True,
            )
            async for chunk in stream:
                delta = chunk.choices[0].delta.content or ""
                if delta:
                    yield delta.encode("utf-8")
            yield b"\n\ndata: [DONE]"
        except Exception as e:
            err = json.dumps({"__stream_error__": str(e)}, ensure_ascii=False)
            yield f"\n\ndata: [ERR]{err}".encode("utf-8")

    return StreamingResponse(
        stream_gen(),
        media_type="text/plain; charset=utf-8",
        headers={"X-Accel-Buffering": "no"},
    )


@app.post("/api/analyze-definition")
async def analyze_definition(req: ExportRequest,
                              x_api_key: Optional[str] = Header(None),
                              x_api_url: Optional[str] = Header(None),
                              x_api_model: Optional[str] = Header(None)):
    client = _get_client(x_api_key, x_api_url)
    model  = _get_model(x_api_model)
    history_text = "\n\n".join(
        f"{'业务方' if m.role == 'user' else '分析师'}：{m.content}"
        for m in req.history
    )
    req_text = "\n".join([
        f"- {r.get('module','')}: {r.get('feature','')} — {r.get('description','')}"
        if isinstance(r, dict) else f"- {r}"
        for r in req.requirements
    ]) if req.requirements else "（暂无结构化需求，请根据访谈记录分析）"

    prompt = f"""你是朝曦金融科技架构顾问，精通 Claude Agent SDK 中 Skill 和 Agent 的设计原则。

请根据以下访谈记录和已梳理需求，判断该需求更适合实现为 Skill 还是 Agent，并给出详细分析。

## 访谈记录
{history_text}

## 已梳理需求
{req_text}

## 判断标准参考

**Skill（专项技能）适合以下情形：**
- 输入输出明确，逻辑相对固定
- 单一职责，完成一件具体的事
- 无需自主决策或动态规划
- 可被其他系统或 Agent 直接调用
- 执行时间短，结果可预期

**Agent（自主智能体）适合以下情形：**
- 需要多步骤推理和动态决策
- 要根据中间结果调整后续行动
- 需要协调调用多个工具或 Skill
- 有较长任务生命周期或需要持久化状态
- 能处理异常情况并自主恢复

请严格按以下 JSON 格式输出，不要输出任何其他文字：
{{
  "recommendation": "Skill" 或 "Agent" 或 "Skill + Agent 组合",
  "confidence": 0到100的整数,
  "summary": "一句话核心结论，不超过40字",
  "skill_fit": 0到100整数,
  "agent_fit": 0到100整数,
  "reasons_for": ["推荐此方案的核心理由1（结合具体需求场景）", "理由2", "理由3"],
  "reasons_against": ["不宜用另一种方式的理由1", "理由2"],
  "skill_roles": ["若涉及Skill，它负责的具体模块1", "模块2"],
  "agent_roles": ["若涉及Agent，它负责的具体能力1", "能力2"],
  "architecture_suggestion": "2-3句话的架构落地建议，结合朝曦业务场景"
}}"""

    try:
        response = client.chat.completions.create(
            model=model, max_tokens=1500,
            messages=[{"role": "user", "content": prompt}]
        )
    except Exception as e:
        raise HTTPException(500, f"分析失败：{e}")

    raw = response.choices[0].message.content
    cleaned = re.sub(r"```(?:json)?\s*", "", raw).replace("```", "").strip()
    try:
        return json.loads(cleaned)
    except Exception:
        match = re.search(r"\{[\s\S]*\}", cleaned)
        if match:
            try:
                return json.loads(match.group())
            except Exception:
                pass
    raise HTTPException(500, "分析结果解析失败，请重试")


@app.post("/api/export/word")
async def export_word(req: WordExportReq):
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    import urllib.parse

    doc = Document()

    style = doc.styles['Normal']
    style.font.name = '微软雅黑'
    style.font.size = Pt(11)

    if req.type == "requirements":
        title = doc.add_heading('需求梳理报告', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph(f'生成时间：{__import__("datetime").datetime.now().strftime("%Y年%m月%d日 %H:%M")}').alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph('')

        for i, r in enumerate(req.requirements):
            if isinstance(r, str):
                doc.add_heading(f'{i+1}. {r}', 2)
                continue
            h = doc.add_heading(f'{i+1}. {r.get("module","功能")} · {r.get("feature","")}', 1)
            p = doc.add_paragraph()
            p.add_run('功能说明：').bold = True
            p.add_run(r.get('description', ''))
            if r.get('process'):
                doc.add_paragraph().add_run('处理流程：').bold = True
                for j, step in enumerate(r['process']):
                    doc.add_paragraph(f'{j+1}. {step}', style='List Number')
            if r.get('inputs'):
                p2 = doc.add_paragraph()
                p2.add_run('输入数据：').bold = True
                p2.add_run('、'.join(r['inputs']))
            if r.get('outputs'):
                p3 = doc.add_paragraph()
                p3.add_run('输出结果：').bold = True
                p3.add_run('、'.join(r['outputs']))
            if r.get('wireframe'):
                p4 = doc.add_paragraph()
                p4.add_run('界面布局：').bold = True
                p4.add_run(r['wireframe'])
            doc.add_paragraph('─' * 40)

        filename = '需求梳理报告.docx'

    elif req.type == "prd":
        lines = req.prd_text.split('\n')
        for line in lines:
            line_stripped = line.strip()
            if not line_stripped:
                doc.add_paragraph('')
                continue
            if line_stripped.startswith('# '):
                h = doc.add_heading(line_stripped[2:], 0)
                h.alignment = WD_ALIGN_PARAGRAPH.CENTER
            elif line_stripped.startswith('## '):
                doc.add_heading(line_stripped[3:], 1)
            elif line_stripped.startswith('### '):
                doc.add_heading(line_stripped[4:], 2)
            elif line_stripped.startswith('#### '):
                doc.add_heading(line_stripped[5:], 3)
            elif line_stripped.startswith(('- ', '* ', '• ')):
                doc.add_paragraph(line_stripped[2:], style='List Bullet')
            elif line_stripped and line_stripped[0].isdigit() and '. ' in line_stripped[:4]:
                doc.add_paragraph(line_stripped, style='List Number')
            elif line_stripped.startswith('**') and line_stripped.endswith('**'):
                p = doc.add_paragraph()
                p.add_run(line_stripped.strip('*')).bold = True
            else:
                p = doc.add_paragraph()
                import re as _re
                parts = _re.split(r'\*\*(.+?)\*\*', line_stripped)
                for k, part in enumerate(parts):
                    run = p.add_run(part)
                    if k % 2 == 1:
                        run.bold = True

        filename = 'PRD需求规格说明书.docx'
    else:
        raise HTTPException(400, "未知导出类型")

    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)

    encoded = urllib.parse.quote(filename)
    return Response(
        content=buffer.getvalue(),
        media_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        headers={'Content-Disposition': f"attachment; filename*=UTF-8''{encoded}"}
    )


@app.post("/api/transcribe")
async def transcribe_audio(file: UploadFile = File(...),
                           x_api_key: Optional[str] = Header(None),
                           x_api_url: Optional[str] = Header(None),
                           x_api_model: Optional[str] = Header(None)):
    """接收录音文件，调用 Whisper 转文字（语音模型固定为 whisper-1）"""
    client = _get_client(x_api_key, x_api_url)
    try:
        audio_bytes = await file.read()
        if not audio_bytes:
            return {"error": "收到空文件，请重新录音"}

        fname = file.filename or "audio.webm"
        content_type = file.content_type or "audio/webm"
        print(f"[transcribe] 收到文件: {fname}, 类型: {content_type}, 大小: {len(audio_bytes)} bytes")

        transcript = client.audio.transcriptions.create(
            model="whisper-1",
            file=(fname, audio_bytes, content_type),
            language="zh"
        )
        text = transcript.text.strip()
        print(f"[transcribe] 识别结果: {text[:80]}")
        return {"text": text}

    except Exception as e:
        err_msg = str(e)
        print(f"[transcribe] 错误: {err_msg}")
        return {"text": "", "error": f"识别失败：{err_msg[:200]}"}


@app.post("/api/generate-demo")
async def generate_demo(req: DemoRequest,
                        x_api_key: Optional[str] = Header(None),
                        x_api_url: Optional[str] = Header(None),
                        x_api_model: Optional[str] = Header(None)):
    """根据需求访谈内容生成可交互的HTML原型Demo"""
    client = _get_client(x_api_key, x_api_url)
    model  = _get_model(x_api_model)
    history_text = "\n".join(
        f"{'用户' if m.role=='user' else 'AI分析师'}: {m.content[:300]}"
        for m in req.history[-30:]
    )
    req_text = ""
    if req.requirements:
        items = req.requirements if isinstance(req.requirements[0], str) else [
            f"- {r.get('feature','')}: {r.get('description','')}" for r in req.requirements
        ]
        req_text = "\n".join(items[:15])

    prompt = f"""你是一位资深前端工程师，擅长快速制作高保真交互原型。

【需求访谈记录】
{history_text}

【整理后的功能需求】
{req_text if req_text else '（请从访谈记录中提取）'}

请根据以上需求，生成一个完整的、可交互的HTML原型Demo，要求：

1. **完全自包含**：单个HTML文件，内联CSS和JavaScript，不依赖任何外部CDN或资源
2. **真实可用**：包含与需求匹配的模拟数据（3-8条示例记录），用户可以真实操作（点击、筛选、填写表单、查看详情等）
3. **专业外观**：金融/企业级UI风格，配色深蓝+白，简洁专业；使用系统字体
4. **核心流程完整**：覆盖需求中最关键的1-2个主流程，让用户能感受到产品的核心价值
5. **导航清晰**：如有多个功能模块，用标签页或侧边栏组织，每个模块都可点击操作
6. **数据互动**：支持增删改查中至少2种操作（如新增记录、筛选查询、查看详情、状态变更等）
7. **朝曦场景**：数据内容贴合家族财富/税务/资本市场/法律/企业治理等专业场景

注意：
- 不要使用 alert()，改用页内提示
- 表格数据要真实，字段名称要专业
- 按钮点击要有反馈（高亮、状态变化、列表刷新等）
- 顶部显示"⚡ 需求原型Demo — [系统名称]（仅供需求确认，非最终产品）"

只输出完整的HTML代码，从<!DOCTYPE html>开始，不要有任何解释文字。"""

    try:
        resp = client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": "你是专业的前端原型工程师，直接输出完整HTML代码，不含任何Markdown标记或解释。"},
                {"role": "user", "content": prompt}
            ],
            max_tokens=6000,
            temperature=0.4
        )
        html = resp.choices[0].message.content.strip()
        if html.startswith("```"):
            html = html.split("```", 2)[1]
            if html.startswith("html"):
                html = html[4:]
            html = html.rsplit("```", 1)[0].strip()
        if "<html" not in html.lower() and "<!doctype" not in html.lower():
            raise HTTPException(500, f"API返回内容不是有效HTML，可能被截断。原始内容前200字：{html[:200]}")
        return {"html": html}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"Demo生成失败：{type(e).__name__}: {str(e)[:300]}")


# ─────────────────────────────────────────────
# 用户认证 & 会话云端存储
# ─────────────────────────────────────────────

@app.post("/api/auth/login")
async def login(req: LoginRequest):
    """邮箱登录（无密码，内部工具）。首次自动注册，老用户直接返回 token。"""
    try:
        email = req.email.strip().lower()
        name  = req.name.strip()
        if not email or "@" not in email:
            raise HTTPException(400, "请输入有效的邮箱地址")
        if not name:
            raise HTTPException(400, "请输入您的姓名")

        admin_email = os.environ.get("ADMIN_EMAIL", "").strip().lower()
        is_admin = 1 if email == admin_email else 0

        print(f"[login] 尝试登录: {email}")
        rows = _sb("GET", "/rest/v1/users", params={"email": f"eq.{email}", "select": "*"})
        now  = datetime.now().isoformat()

        if rows and isinstance(rows, list):
            user = rows[0]
            new_admin = max(is_admin, int(user.get("is_admin", 0)))
            _sb("PATCH", "/rest/v1/users",
                data={"name": name, "last_active": now, "is_admin": new_admin},
                params={"email": f"eq.{email}"})
            print(f"[login] 老用户登录成功: {email}")
            return {"token": user["token"], "user_id": user["id"],
                    "user": {"name": name, "email": email, "is_admin": bool(new_admin)}}
        else:
            uid   = secrets.token_hex(8)
            token = secrets.token_hex(24)
            _sb("POST", "/rest/v1/users", data={
                "id": uid, "email": email, "name": name, "token": token,
                "is_admin": is_admin, "created_at": now, "last_active": now
            })
            print(f"[login] 新用户注册成功: {email}")
            return {"token": token, "user_id": uid,
                    "user": {"name": name, "email": email, "is_admin": bool(is_admin)}}
    except HTTPException:
        raise
    except Exception as e:
        print(f"[login] 未捕获异常: {type(e).__name__}: {e}")
        raise HTTPException(500, f"登录异常: {type(e).__name__}: {str(e)[:200]}")


@app.get("/api/auth/me")
async def get_me(authorization: Optional[str] = Header(None)):
    user = require_token(authorization)
    return {"user_id": user["id"], "name": user["name"],
            "email": user["email"], "is_admin": bool(user.get("is_admin"))}


@app.post("/api/sessions/save")
async def save_session(req: SaveSessionRequest, authorization: Optional[str] = Header(None)):
    """保存或更新一条会话到服务端"""
    user = require_token(authorization)
    now  = datetime.now().isoformat()
    sid  = req.session_id or secrets.token_hex(10)

    h_json = json.dumps(req.history,      ensure_ascii=False)
    r_json = json.dumps(req.requirements, ensure_ascii=False)
    # 将附加字段打包进 mermaid 列（JSON envelope，避免 schema 变更）
    if req.wireframe_mermaid or req.prd_content or req.ui_wireframe_mermaid:
        mermaid_data = json.dumps({
            "v": 1,
            "flowchart": req.mermaid,
            "wireframe": req.wireframe_mermaid,
            "ui_wireframe": req.ui_wireframe_mermaid,
            "prd": req.prd_content
        }, ensure_ascii=False)
    else:
        mermaid_data = req.mermaid

    existing = _sb("GET", "/rest/v1/sessions",
                   params={"id": f"eq.{sid}", "user_id": f"eq.{user['id']}", "select": "id"})

    if existing and isinstance(existing, list):
        _sb("PATCH", "/rest/v1/sessions",
            data={"history": h_json, "requirements": r_json, "mermaid": mermaid_data,
                  "completeness": req.completeness, "title": req.title, "updated_at": now},
            params={"id": f"eq.{sid}"})
    else:
        _sb("POST", "/rest/v1/sessions", data={
            "id": sid, "user_id": user["id"], "title": req.title,
            "history": h_json, "requirements": r_json, "mermaid": mermaid_data,
            "completeness": req.completeness, "created_at": now, "updated_at": now
        })
    return {"session_id": sid, "updated_at": now}


@app.get("/api/sessions")
async def list_sessions(authorization: Optional[str] = Header(None)):
    """获取当前用户的所有会话列表（不含完整 history）"""
    user = require_token(authorization)
    rows = _sb("GET", "/rest/v1/sessions",
               params={"user_id": f"eq.{user['id']}",
                       "select": "id,title,completeness,created_at,updated_at",
                       "order": "updated_at.desc"})
    return {"sessions": rows if isinstance(rows, list) else []}


@app.get("/api/sessions/{session_id}")
async def get_session(session_id: str, authorization: Optional[str] = Header(None)):
    """获取单条会话完整内容"""
    user = require_token(authorization)
    rows = _sb("GET", "/rest/v1/sessions",
               params={"id": f"eq.{session_id}", "user_id": f"eq.{user['id']}", "select": "*"})
    if not rows or not isinstance(rows, list):
        raise HTTPException(404, "会话不存在")
    d = dict(rows[0])
    d["history"]      = json.loads(d.get("history")      or "[]")
    d["requirements"] = json.loads(d.get("requirements") or "[]")
    # 解包 mermaid 列中可能存在的 JSON 包（含 wireframeMermaid + prdContent）
    raw_mermaid = d.get("mermaid") or ""
    if raw_mermaid.startswith('{"v":1'):
        try:
            m = json.loads(raw_mermaid)
            d["mermaid"]            = m.get("flowchart", "")
            d["wireframeMermaid"]   = m.get("wireframe", "")
            d["uiWireframeMermaid"] = m.get("ui_wireframe", "")
            d["prdContent"]         = m.get("prd", "")
        except Exception:
            d["wireframeMermaid"]   = ""
            d["uiWireframeMermaid"] = ""
            d["prdContent"]         = ""
    else:
        d["wireframeMermaid"]   = ""
        d["uiWireframeMermaid"] = ""
        d["prdContent"]         = ""
    return d


@app.delete("/api/sessions/{session_id}")
async def delete_session(session_id: str, authorization: Optional[str] = Header(None)):
    user = require_token(authorization)
    _sb("DELETE", "/rest/v1/sessions",
        params={"id": f"eq.{session_id}", "user_id": f"eq.{user['id']}"})
    return {"ok": True}


# ─────────────────────────────────────────────
# 管理员看板
# ─────────────────────────────────────────────

@app.get("/api/admin/stats")
async def admin_stats(authorization: Optional[str] = Header(None)):
    require_admin(authorization)
    total_users    = _sb_count("/rest/v1/users")
    total_sessions = _sb_count("/rest/v1/sessions")
    seven_days_ago = (datetime.now() - timedelta(days=7)).isoformat()
    active_rows    = _sb("GET", "/rest/v1/sessions",
                         params={"updated_at": f"gte.{seven_days_ago}", "select": "user_id"})
    active = len(set(r["user_id"] for r in (active_rows if isinstance(active_rows, list) else [])))
    return {"total_users": total_users, "total_sessions": total_sessions, "active_7d": active}


@app.get("/api/admin/users")
async def admin_users(authorization: Optional[str] = Header(None)):
    require_admin(authorization)
    rows = _sb("GET", "/rest/v1/users",
               params={"is_admin": "eq.0",
                       "select": "id,name,email,created_at,last_active,sessions(id,completeness,updated_at)",
                       "order": "last_active.desc.nullslast"})
    result = []
    for u in (rows if isinstance(rows, list) else []):
        sessions_data = u.pop("sessions", None) or []
        u["session_count"]    = len(sessions_data)
        u["last_session"]     = max(
            (s["updated_at"] for s in sessions_data if s.get("updated_at")), default=None)
        u["max_completeness"] = max(
            (s.get("completeness") or 0 for s in sessions_data), default=0)
        result.append(u)
    return {"users": result}


@app.get("/api/admin/users/{user_id}/sessions")
async def admin_user_sessions(user_id: str, authorization: Optional[str] = Header(None)):
    require_admin(authorization)
    rows = _sb("GET", "/rest/v1/sessions",
               params={"user_id": f"eq.{user_id}",
                       "select": "id,title,completeness,created_at,updated_at",
                       "order": "updated_at.desc"})
    return {"sessions": rows if isinstance(rows, list) else []}


@app.get("/api/admin/sessions/{session_id}")
async def admin_get_session(session_id: str, authorization: Optional[str] = Header(None)):
    require_admin(authorization)
    rows = _sb("GET", "/rest/v1/sessions",
               params={"id": f"eq.{session_id}", "select": "*"})
    if not rows or not isinstance(rows, list):
        raise HTTPException(404, "会话不存在")
    d = dict(rows[0])
    d["history"]      = json.loads(d.get("history")      or "[]")
    d["requirements"] = json.loads(d.get("requirements") or "[]")
    # 解包 mermaid 列中可能存在的 JSON 包（含 wireframeMermaid + prdContent）
    raw_mermaid = d.get("mermaid") or ""
    if raw_mermaid.startswith('{"v":1'):
        try:
            m = json.loads(raw_mermaid)
            d["mermaid"]            = m.get("flowchart", "")
            d["wireframeMermaid"]   = m.get("wireframe", "")
            d["uiWireframeMermaid"] = m.get("ui_wireframe", "")
            d["prdContent"]         = m.get("prd", "")
        except Exception:
            d["wireframeMermaid"]   = ""
            d["uiWireframeMermaid"] = ""
            d["prdContent"]         = ""
    else:
        d["wireframeMermaid"]   = ""
        d["uiWireframeMermaid"] = ""
        d["prdContent"]         = ""
    return d


@app.get("/admin", response_class=HTMLResponse)
async def admin_page():
    """管理员看板页面"""
    content = open("static/admin.html", encoding="utf-8").read() \
              if os.path.exists("static/admin.html") else "<h1>管理页面未找到</h1>"
    return HTMLResponse(content, headers={
        "Cache-Control": "no-cache, no-store, must-revalidate",
        "Pragma": "no-cache", "Expires": "0"
    })

@app.get("/", response_class=HTMLResponse)
async def index_page():
    """主页面 — 禁止浏览器缓存，确保每次刷新都取最新 HTML"""
    content = open("static/index.html", encoding="utf-8").read()
    return HTMLResponse(content, headers={
        "Cache-Control": "no-cache, no-store, must-revalidate",
        "Pragma": "no-cache", "Expires": "0"
    })

app.mount("/", StaticFiles(directory="static", html=True), name="static")

if __name__ == "__main__":
    import uvicorn
    port = int(os.environ.get("PORT", 8000))
    uvicorn.run(app, host="0.0.0.0", port=port, reload=False)
